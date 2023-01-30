import collections 
import collections.abc
import json
import os
import traceback
from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.chart.axis import _BaseAxis
from pptx.enum.chart import XL_CATEGORY_TYPE
from pptx.dml.color import ColorFormat
from pptx.chart.series import _BaseSeries
from pptx.dml.chtfmt import ChartFormat
from pptx.dml.fill import FillFormat
from pptx.enum.chart import XL_TICK_MARK
from pptx.enum.chart import XL_TICK_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.series import AreaSeries
from pptx.enum.chart import XL_MARKER_STYLE

# def ukuran(data):
#         try:
#             return round((data / 96),2)
#         except Exception:
#             traceback.print_exc()

# with open(R"src\charts\config\area_config.json","r") as f:
#     data = json.load(f)
    
# with open(R"src\charts\config\area_config.json","r") as f:
#     data = json.load(f)
    
# prs = Presentation()
# prs.slide_width = Inches(13.333)
# prs.slide_height = Inches(7.5)


# slides = data.get('slides')


prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Tambahkan chart pada slide
chart_data = ChartData()
chart_data.categories = ['pks', 'gerindra', 'psi', 'pdi-p', 'ppp', 'pbb', 'demokrat', 'pkb', 'golkar',  'pan']
chart_data.add_series('2022', [336, 196, 161, 55, 35, 31, 31, 27, 20, 16])
chart_data.add_series('2023', [12, 16, 61, 55, 35, 38, 39, 97, 20, 6])

x, y, cx, cy = Inches(6.87), Inches(1.59), Inches(5.93), Inches(4.33)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
).chart


data_label = chart.plots[0]
data_label.has_data_labels = True
data_labels = data_label.data_labels
data_labels.show_value = True
# data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
data_labels.font.size = Pt(11)
chart.has_title = False
data_label.overlap = -15
data_label.gap_width = 54


value_axis = chart.value_axis
value_axis.visible = True
value_axis.has_major_gridlines = False
value_axis.major_tick_mark = XL_TICK_MARK.NONE
value_axis.tick_labels.number_format = 'none'
value_axis.tick_labels.font.size = Pt(10)
value_axis.tick_labels.font.color.rgb = RGBColor.from_string("000000")
# value_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
value_axis.reverse_order = False

# chart.has_category_axis = False 
category_axis = chart.category_axis
category_axis.reverse_order = False
category_axis.major_tick_mark = XL_TICK_MARK.NONE
category_axis.tick_labels.font.size = Pt(10)
# category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW



#Legend =============================================
# chart.has_legend = True
# chart.legend.include_in_layout = True
# chart.legend.horz_offset = -1.0
# chart.legend.font.bold = True
# chart.legend.font.size = Pt(50)
# chart.legend.font.color.rgb = RGBColor(191,191,191)
# chart.legend.font.underline = True
# chart.legend.position = XL_LEGEND_POSITION.BOTTOM



#Series =============================================
series = chart.plots[0].series[0]
series.smooth = True
series.format.line.fill.solid()
series.format.line.color.rgb = RGBColor.from_string('B4D92A')
series.marker.style = XL_MARKER_STYLE.CIRCLE

series = chart.plots[0].series[1]
series.smooth = True
series.format.fill.solid()
series.format.line.color.rgb = RGBColor.from_string('1DB7D9')
series.marker.style = XL_MARKER_STYLE.CIRCLE


textbox = slide.shapes.add_textbox(Inches(1.05),   #left
                                    Inches(1.05), #top
                                    Inches(4.44), #width
                                    Inches(0.4)) #height
                                    
tf = textbox.text_frame
# tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.font.size = Pt(18)

p.alignment = PP_ALIGN.CENTER

run.text = 'LINE'


#CHART 2
chart_data = ChartData()
chart_data.categories = ['pks', 'gerindra', 'psi', 'pdi-p', 'ppp', 'pbb', 'demokrat', 'pkb', 'golkar',  'pan']
chart_data.add_series('', [336, 196, 161, 55, 35, 31, 31, 27, 20, 16])
chart_data.add_series('', [12, 16, 61, 55, 35, 38, 39, 97, 20, 6])

x, y, cx, cy = Inches(0.35), Inches(1.59), Inches(5.93), Inches(4.33)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
).chart

textbox = slide.shapes.add_textbox(Inches(7.84),   #left
                                    Inches(1.05), #top
                                    Inches(4.44), #width
                                    Inches(0.4)) #height
tf = textbox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.font.size = Pt(18)

p.alignment = PP_ALIGN.CENTER

run.text = 'LINE MODIFIED'







#========================================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Tambahkan chart pada slide
chart_data = ChartData()
chart_data.categories = ['pks', 'gerindra', 'psi', 'pdi-p', 'ppp', 'pbb', 'demokrat', 'pkb', 'golkar',  'pan']
chart_data.add_series('2022', [336, 196, 161, 55, 35, 31, 31, 27, 20, 16])
chart_data.add_series('2023', [12, 16, 61, 55, 35, 38, 39, 97, 20, 6])

x, y, cx, cy = Inches(6.87), Inches(1.59), Inches(5.93), Inches(4.33)
chart = slide2.shapes.add_chart(
    XL_CHART_TYPE.LINE_MARKERS_STACKED, x, y, cx, cy, chart_data
).chart


data_label = chart.plots[0]
data_label.has_data_labels = True
data_labels = data_label.data_labels
data_labels.show_value = True
# data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
data_labels.font.size = Pt(11)
chart.has_title = False
data_label.overlap = -15
data_label.gap_width = 54


value_axis = chart.value_axis
value_axis.visible = True
value_axis.has_major_gridlines = False
value_axis.major_tick_mark = XL_TICK_MARK.NONE
value_axis.tick_labels.number_format = 'none'
value_axis.tick_labels.font.size = Pt(10)
value_axis.tick_labels.font.color.rgb = RGBColor.from_string("000000")
# value_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
value_axis.reverse_order = False

# chart.has_category_axis = False 
category_axis = chart.category_axis
category_axis.reverse_order = False
category_axis.major_tick_mark = XL_TICK_MARK.NONE
category_axis.tick_labels.font.size = Pt(10)
# category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW



#Legend =============================================
# chart.has_legend = True
# chart.legend.include_in_layout = True
# chart.legend.horz_offset = -1.0
# chart.legend.font.bold = True
# chart.legend.font.size = Pt(50)
# chart.legend.font.color.rgb = RGBColor(191,191,191)
# chart.legend.font.underline = True
# chart.legend.position = XL_LEGEND_POSITION.BOTTOM



#Series =============================================
series = chart.plots[0].series[0]
series.smooth = True
series.format.line.fill.solid()
series.format.line.color.rgb = RGBColor.from_string('B4D92A')
series.marker.style = XL_MARKER_STYLE.CIRCLE

series = chart.plots[0].series[1]
series.smooth = True
series.format.fill.solid()
series.format.line.color.rgb = RGBColor.from_string('1DB7D9')
series.marker.style = XL_MARKER_STYLE.CIRCLE

textbox = slide2.shapes.add_textbox(Inches(1.05),   #left
                                    Inches(1.05), #top
                                    Inches(4.44), #width
                                    Inches(0.4)) #height
                                    
tf = textbox.text_frame
# tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.font.size = Pt(18)

p.alignment = PP_ALIGN.CENTER

run.text = 'LINE MARKER STACKED'


#CHART 2
chart_data = ChartData()
chart_data.categories = ['pks', 'gerindra', 'psi', 'pdi-p', 'ppp', 'pbb', 'demokrat', 'pkb', 'golkar',  'pan']
chart_data.add_series('', [336, 196, 161, 55, 35, 31, 31, 27, 20, 16])
chart_data.add_series('', [12, 16, 61, 55, 35, 38, 39, 97, 20, 6])

x, y, cx, cy = Inches(0.35), Inches(1.59), Inches(5.93), Inches(4.33)
chart = slide2.shapes.add_chart(
    XL_CHART_TYPE.LINE_MARKERS_STACKED, x, y, cx, cy, chart_data
).chart

textbox = slide2.shapes.add_textbox(Inches(7.84),   #left
                                    Inches(1.05), #top
                                    Inches(4.44), #width
                                    Inches(0.4)) #height
tf = textbox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.font.size = Pt(18)

p.alignment = PP_ALIGN.CENTER

run.text = 'LINE MARKERS STACKED MODIFIED'







# Simpan presentasi
prs.save(R'src\charts\result\line.pptx')                        
os.startfile(R'src\charts\result\line.pptx')
