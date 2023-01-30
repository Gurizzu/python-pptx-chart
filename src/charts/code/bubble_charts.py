import collections 
import collections.abc
import json
import os
import traceback
from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.chart.axis import _BaseAxis
from pptx.enum.chart import XL_CATEGORY_TYPE
from pptx.dml.color import ColorFormat
from pptx.chart.series import _BaseSeries
from pptx.dml.chtfmt import ChartFormat
from pptx.dml.fill import FillFormat
from pptx.enum.chart import XL_TICK_MARK
from pptx.enum.chart import XL_TICK_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import BubbleChartData


prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])

chart_data = BubbleChartData()

series_1 = chart_data.add_series('Series 1')
series_1.add_data_point(0.7, 2.7, 10)
series_1.add_data_point(1.8, 3.2, 4)
series_1.add_data_point(2.6, 0.8, 8)

x, y, cx, cy = Inches(0.35), Inches(1.59), Inches(4.51), Inches(3.1)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.BUBBLE, x, y, cx, cy, chart_data
).chart

# value_axis = chart.value_axis
# value_axis.visible = True
# value_axis.has_major_gridlines = False
# value_axis.major_tick_mark = XL_TICK_MARK.NONE
# value_axis.tick_labels.number_format = 'none'
# value_axis.tick_labels.font.size = Pt(10)
# value_axis.tick_labels.font.color.rgb = RGBColor.from_string("000000")
# value_axis.tick_label_position = XL_TICK_LABEL_POSITION.HIGH
# value_axis.reverse_order = False

# # chart.has_category_axis = False 
# category_axis = chart.category_axis
# category_axis.reverse_order = True
# category_axis.major_tick_mark = XL_TICK_MARK.NONE
# category_axis.tick_labels.font.size = Pt(10)
# # category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW

chart_series = chart.series[0]
chart_series.format.fill.solid()
chart_series.format.fill.fore_color.rgb = RGBColor.from_string('B4D92A')

# chart_series = chart.series[1]
# chart_series.format.fill.solid()
# chart_series.format.fill.fore_color.rgb = RGBColor.from_string('1DB7D9')


prs.save(R'src\charts\result\bubble.pptx')
os.startfile(R'src\charts\result\bubble.pptx')
