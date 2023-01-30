import collections 
import collections.abc
import json
import os
import traceback
from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.chart.axis import _BaseAxis
from pptx.enum.chart import XL_CATEGORY_TYPE
from pptx.dml.color import ColorFormat
from pptx.chart.series import _BaseSeries
from pptx.dml.chtfmt import ChartFormat
from pptx.dml.fill import FillFormat
from pptx.enum.chart import XL_TICK_MARK
from pptx.enum.chart import XL_TICK_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import BubbleChartData


prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])

chart_data = ChartData()
chart_data.categories = ['West', 'East', 'North', 'South', 'Other']
chart_data.add_series('Series 1', (0.135, 0.324, 0.180, 0.235, 0.126))


x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
).chart

chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False

chart.plots[0].has_data_labels = True
data_labels = chart.plots[0].data_labels
data_labels.number_format = '0%'
data_labels.position = XL_LABEL_POSITION.OUTSIDE_END


prs.save(R'src\charts\result\pie.pptx')
os.startfile(R'src\charts\result\pie.pptx')
