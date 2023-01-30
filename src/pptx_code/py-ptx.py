import collections 
import collections.abc
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.util import Inches, Pt


prs = Presentation('template_1.pptx')

slide = prs.slides.add_slide(prs.slide_layouts[0])
for shape in slide.placeholders:
    print('%d %s' % (shape.placeholder_format.idx, shape.name))
    
date_shape = slide.placeholders[14]
title = slide.placeholders[0]
pychart = slide.placeholders[13]
date_shape.text = "criteria_name"
title.text = "ini judul"

# table_shape = slide.placeholders[13]
# table = table_shape.insert_table(rows=10, cols=2).table

chart_data = CategoryChartData()
chart_data.categories = ["positive", "neutral", "negative"]
chart_data.add_series("Tes", (100,200,90))

# add chart to slide --------------------
x, y, cx, cy = Inches(8.08), Inches(1.00), Inches(1.73), Inches(1.60)
shapes = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data)
chart = shapes.chart
chart.font.size = Pt(8)
chart.font.color.rgb = RGBColor(255, 255, 255)
chart.plots[0].has_data_labels = True
data_labels = chart.plots[0].data_labels
data_labels.number_format = '0%'
data_labels.position = XL_LABEL_POSITION.INSIDE_END

color_list = ["007fff", "696969", "8B0000"]
pychart.insert_chart(XL_CHART_TYPE.DOUGHNUT, chart_data)

for idx, point in enumerate(slide.placeholders[13].chart.series[0].points):
                col_idx = idx % len(color_list)
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = RGBColor.from_string(color_list[col_idx])


prs.save('test.pptx')