import collections 
import collections.abc
import json
import os
from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.chart.axis import _BaseAxis
from pptx.enum.chart import XL_CATEGORY_TYPE
from pptx.dml.color import ColorFormat
from pptx.chart.series import _BaseSeries
from pptx.dml.chtfmt import ChartFormat
from pptx.dml.fill import FillFormat
from pptx.enum.chart import XL_TICK_MARK
from pptx.enum.chart import XL_TICK_LABEL_POSITION
import traceback

from utils import chart_type_check, text_type

def ukuran(data):
        try:
            return round((data / 96),2)
        except Exception:
            traceback.print_exc()

with open("config.json","r") as f:
    data = json.load(f)

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slides = data.get('slides')
for slide in slides:
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    for widget in slide.get("widgets"):
        match widget.get("widget").get("type"):
            case "chart":
                chart_data = ChartData()
                chart_data.categories = widget.get("widget").get("widget_data").get("categories")
                for chart_series in widget.get("widget").get("widget_data").get("series"):
                    chart_data.add_series(chart_series.get("legend"),chart_series.get("point"))
                    
                x, y, cx, cy = Inches(widget.get("left")), Inches(widget.get("top")), Inches(widget.get("height")), Inches(widget.get("width"))
                
                chart_type = widget.get("widget").get("widget_ops").get("type")
                chart = chart_type_check(chart_type, slide1, chart_data, x, y, cx, cy)
                
                # widget_format = widget.get("widget").get("widget_ops").get("format")
                
                # value_axis = chart.value_axis
                # if widget_format.get("axis").get("value").get("has_text_frame") == True:
                #     value_axis.has_major_gridlines = True
                # else:
                #     value_axis.has_major_gridlines = False
                    
                # match widget_format.get("axis").get("value").get("major_tick_mark"):
                #     case "NONE":
                #         value_axis.major_tick_mark = XL_TICK_MARK.NONE
                #     case "INSIDE":
                #         value_axis.major_tick_mark = XL_TICK_MARK.INSIDE
                #     case "OUTSIDE":
                #         value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE
                        
                # number_format = widget_format.get("axis").get("value").get("number_format")
                # value_axis.tick_labels.number_format = f'{number_format}'
                    
                
                
            
            case "text":
                # x, y, cx, cy = Inches(widget.get("left")), Inches(widget.get("top")), Inches(widget.get("height")), Inches(widget.get("width"))
                textbox = slide1.shapes.add_textbox(Inches(widget.get("left")),
                                                    Inches(widget.get("top")) ,
                                                    width = Inches(widget.get("width")),
                                                    height = Inches(widget.get("height")))
                tf = textbox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                run = p.add_run()
                run.font.size = Pt(18)
                p.alignment = PP_ALIGN.JUSTIFY
                run.text = widget.get("widget").get("textbox").get("text")
                
                        
prs.save('chart_example3.pptx')                        
                    
                

# prs.save('chart_example_2.pptx')
# os.startfile("chart_example_2.pptx")