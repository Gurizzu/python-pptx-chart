import collections 
import collections.abc
import json
import os
from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.chart.axis import _BaseAxis
from pptx.enum.chart import XL_CATEGORY_TYPE
from pptx.dml.color import ColorFormat
from pptx.chart.series import _BaseSeries
from pptx.dml.chtfmt import ChartFormat
from pptx.dml.fill import FillFormat
from pptx.enum.chart import XL_TICK_MARK
from pptx.enum.chart import XL_TICK_LABEL_POSITION
import traceback

def chart_type_check(chart_type:str,slides,chart_data,x, y, cx, cy):
    # print(chart_type)
    match chart_type:
        case "AREA":
            chart = slides.shapes.add_chart(
                XL_CHART_TYPE.AREA, x, y, cx, cy, chart_data
            ).chart
            
        
        case "PIE":
            chart = slides.shapes.add_chart(
                XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
            ).chart
            
        
            
    return chart
            
    
def text_type(slides,text,x, y, cx, cy):
    textbox = slides.shapes.add_textbox(Inches(x),
                                        Inches(y) ,
                                        Inches(cx),
                                        Inches(cy))
    
    return textbox
    
    # tf = textbox.text_frame
    # tf.word_wrap = True
    # p = tf.paragraphs[0]
    # run = p.add_run()
    # run.font.size = Pt(18)
    # p.alignment = PP_ALIGN.JUSTIFY
    # run.text = text
                
                
                
                
                
                
def contoh1(data:dict): 
     
    if data:
        
        print("data masuk")
        
        if len(data) > 5:
            print("data lebih dari 5")
        else:
            print("data kurang dari 5")
            
    else:
        print("data kosong / tidak masuk")
        

def contoh2(data:dict):
    
    if not data:
         print("data kosong / tidak masuk")        
    print("data masuk")
    
    if len(data) > 5:
            print("data lebih dari 5")
    print("data kurang dari 5")
    
    
        
        