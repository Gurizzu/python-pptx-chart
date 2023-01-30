import collections 
import collections.abc
import os
from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.chart.axis import _BaseAxis
from pptx.enum.chart import XL_CATEGORY_TYPE
from pptx.dml.color import ColorFormat
from pptx.chart.series import _BaseSeries
from pptx.dml.chtfmt import ChartFormat
from pptx.dml.fill import FillFormat
from pptx.enum.chart import XL_TICK_MARK
from pptx.enum.chart import XL_TICK_LABEL_POSITION
from pptx.chart.plot import BarPlot
from pptx.chart.datalabel import DataLabels



prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Tambahkan chart pada slide
chart_data = ChartData()
chart_data.categories = ['pks', 'gerindra', 'psi', 'pdi-p', 'ppp', 'pbb', 'demokrat', 'pkb', 'golkar',  'pan']
chart_data.add_series('', [336, 196, 161, 55, 35, 31, 31, 27, 20, 16])

x, y, cx, cy = Inches(5.29), Inches(1.92), Inches(3.74), Inches(4.33)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.BUBBLE, x, y, cx, cy, chart_data
).chart


data_label = chart.plots[0]
data_label.has_data_labels = True
data_labels = data_label.data_labels
data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
data_labels.font.size = Pt(11)

chart.has_title = False

value_axis = chart.value_axis
value_axis.visible = False
value_axis.has_major_gridlines = False
value_axis.major_tick_mark = XL_TICK_MARK.NONE
value_axis.tick_labels.number_format = 'none'
value_axis.tick_labels.font.size = Pt(10)
value_axis.tick_labels.font.color.rgb = RGBColor.from_string("000000")
value_axis.tick_label_position = XL_TICK_LABEL_POSITION.HIGH
value_axis.reverse_order = False

bar_plot = chart.plots[0]
bar_plot.gap_width = 83

# chart.has_category_axis = False 
category_axis = chart.category_axis
category_axis.reverse_order = True
category_axis.major_tick_mark = XL_TICK_MARK.NONE
category_axis.tick_labels.font.size = Pt(10)
# category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW



#Legend =============================================
chart.has_legend = False
# chart.legend.include_in_layout = True
# chart.legend.horz_offset = -1.0
# chart.legend.font.bold = True
# chart.legend.font.size = Pt(50)
# chart.legend.font.color.rgb = RGBColor(191,191,191)
# chart.legend.font.underline = True
# chart.legend.position = XL_LEGEND_POSITION.BOTTOM



#Series =============================================
chart_series = chart.series[0]
chart_series.format.fill.solid()
chart_series.format.fill.fore_color.rgb = RGBColor.from_string('D84E2E')
# fill = chart_series.format.fill
# fill.solid()
# .fore_color.rgb = RGBColor(191,141,191)



# Simpan presentasi
prs.save('bubble.pptx')
os.startfile("bubble.pptx")