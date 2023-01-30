import collections 
import collections.abc
import os
from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR


prs = Presentation('tmplate.pptx')

bulk = {
  "title": "Polda Aceh",
  "date": "3 Agustus - 3 November 2022",
  "OrganizationName": "Polda Aceh",
  "CountNewsOrganization": 1265,
  "Organization_Pie_Positive_News": 646,
  "Organization_Pie_Negative_News": 197,
  "Organization_Pie_Netral_News": 422,
  "Organization_Issue_Positive_News_1": "Kecelakaan",
  "Organization_LinkIssue_Positive_News_1": "https",
  "Organization_Issue_Positive_News_2": "hotline 110",
  "Organization_LinkIssue_Positive_News_2": "https",
  "Organization_Issue_Positive_News_3": "pemberantasan pinjol ilegal",
  "Organization_LinkIssue_Positive_News_3": "https",
  "Organization_Issue_Negative_News_1": "oknum polisi - cyber crime",
  "Organization_LinkIssue_Negative_News_1": "https",
  "Organization_Issue_Negative_News_2": "oknum polisi kekerasan seksual",
  "Organization_LinkIssue_Negative_News_2": "https",
  "Organization_Issue_Negative_News_3": "oknum polisi narkoba",
  "Organization_LinkIssue_Negative_News_3": "https",
  "CountSocmedOrganization": 1265,
  "Organization_Pie_Positive_Socmed": 646,
  "Organization_Pie_Negative_Socmed": 197,
  "Organization_Pie_Netral_Socmed": 422,
  "Organization_OrganicPost": 123,
  "Organization_BotPost": 843,
  "Organization_OrganicAccount": 224,
  "Organization_Bot": 923,
  "Organization_Engagement": 1232,
  "Organization_Issue_Positive_Socmed_1": "kecelakaan",
  "Organization_LinkIssue_Positive_Socmed_1": "https",
  "Organization_Issue_Positive_Socmed_2": "hotline 110",
  "Organization_LinkIssue_Positive_Socmed_2": "https",
  "Organization_Issue_Positive_Socmed_3": "pemberantasan pinjol ilegal",
  "Organization_LinkIssue_Positive_Socmed_3": "https",
  "Organization_Issue_Negative_Socmed_1": "oknum polisi cyber crime",
  "Organization_LinkIssue_Negative_Socmed_1": "https",
  "Organization_Issue_Negative_Socmed_2": "oknum polisi narkoba",
  "Organization_LinkIssue_Negative_Socmed_2": "https",
  "Organization_Issue_Negative_Socmed_3": "oknum polisi pungli",
  "Organization_LinkIssue_Negative_Socmed_3": "https",
  "PersonName": "ngurah rai",
  "CountNewsPerson": 121,
  "Person_Pie_Positive_News": 122,
  "Person_Pie_Negative_News": 123,
  "Person_Pie_Netral_News": 157,
  "Person_Issue_Positive_News_1": "hotline 110",
  "Person_LinkIssue_Positive_News_1": "https",
  "Person_Issue_Positive_News_2": "pemberantasan pinjol",
  "Person_LinkIssue_Positive_News_2": "https",
  "Person_Issue_Positive_News_3": "kecelakaan",
  "Person_LinkIssue_Positive_News_3": "https",
  "Person_Issue_Negative_News_1": "oknum polisi cyber crime",
  "Person_LinkIssue_Negative_News_1": "https",
  "Person_Issue_Negative_News_2": "oknum polisi narkob",
  "Person_LinkIssue_Negative_News_2": "https",
  "Person_Issue_Negative_News_3": "https",
  "Person_LinkIssue_Negative_News_3": "oknum polisi pungli",
  "Person_Pie_Positive_Socmed": "https",
  "Person_Pie_Negative_Socmed": 212,
  "Person_Pie_Netral_Socmed": 23,
  "CountSocmedPerson": 500,
  "Person_OrganicPost": 600,
  "Person_BotPost": 291,
  "Person_OrganicAccount": 473,
  "Person_Bot": 292,
  "Person_Engagement": 8829,
  "Person_Issue_Positive_Socmed_1": "hotline 110",
  "Person_LinkIssue_Positive_Socmed_1": "https",
  "Person_Issue_Positive_Socmed_2": "pemberantasan pinjol",
  "Person_LinkIssue_Positive_Socmed_2": "https",
  "Person_Issue_Positive_Socmed_3": "kecelkaan",
  "Person_LinkIssue_Positive_Socmed_3": "https",
  "Person_Issue_Negative_Socmed_1": "oknum polisi cyber crime",
  "Person_LinkIssue_Negative_Socmed_1": "https",
  "Person_Issue_Negative_Socmed_2": "oknum polisi narkoboy",
  "Person_LinkIssue_Negative_Socmed_2": "https",
  "Person_Issue_Negative_Socmed_3": "oknum polisi pulang",
  "Person_LinkIssue_Negative_Socmed_3": "https"
}



print(len(bulk))

def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell


slide = prs.slides.add_slide(prs.slide_layouts[1])
# for shape in slide.placeholders:
#     print('%d %s' % (shape.placeholder_format.idx, shape.name))
    
slide.placeholders[0].text = bulk.get("title")
slide.placeholders[13].text = bulk.get("date")
slide.placeholders[16].text = str(bulk.get("CountNewsOrganization"))
slide.placeholders[18].text = str(bulk.get("CountSocmedOrganization"))
slide.placeholders[21].text =  bulk.get("date")


#==============================================CHART 1 =============================================

chart_data = ChartData()
data = [197, 646, 422]
categories = ["Negatif", "Positif", "Netral"]
chart_data.categories = categories
chart_data.add_series('Series 1', data)

color_list = ["FF2B2B","3C4DBA","707070"] 

chart2 = slide.placeholders[14]
graphic_frame = chart2.insert_chart(XL_CHART_TYPE.DOUGHNUT, chart_data)
chart = graphic_frame.chart
chart.font.size = Pt(8)
chart.font.color.rgb = RGBColor(255,255,255)
chart.font.bold = True

for idx, point in enumerate(chart.series[0].points):
                col_idx = idx % len(color_list)
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = RGBColor.from_string(color_list[col_idx])

chart.has_legend = False
chart.has_title = False
plot = chart.plots[0]
for i in range(len(data)):
    plot.series[0].points[i].data_label.font.size = Pt(1)
    plot.series[0].points[i].data_label.text_frame.text = f"{categories[i]} \n {str(data[i])}"
    
    
#==============================================CHART 2 =============================================
chart_data = ChartData()
data = [197, 646, 422]
categories = ["Negatif", "Positif", "Netral"]
chart_data.categories = categories
chart_data.add_series('Series 1', data)

color_list = ["FF2B2B","3C4DBA","707070"] 

chart = slide.placeholders[15]
graphic_frame = chart.insert_chart(XL_CHART_TYPE.DOUGHNUT, chart_data)
chart = graphic_frame.chart
chart.font.size = Pt(8)
chart.font.color.rgb = RGBColor(255,255,255)
chart.font.bold = True

for idx, point in enumerate(chart.series[0].points):
                col_idx = idx % len(color_list)
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = RGBColor.from_string(color_list[col_idx])

chart.has_legend = False
chart.has_title = False
plot = chart.plots[0]
for i in range(len(data)):
    plot.series[0].points[i].data_label.font.size = Pt(1)
    plot.series[0].points[i].data_label.text_frame.text = f"{categories[i]} \n {str(data[i])}"



#============================ TABLE ATAS 1 =============================
table = slide.placeholders[17]
x, y, cx, cy = Inches(2), Inches(2), Inches(4), Inches(1.5)
table_shape = table.insert_table(rows=4, cols=2).table
table_shape.columns[0].width = Inches(2.6)
table_shape.columns[1].width = Inches(2.6)
for i in range(4):
    table_shape.rows[i].height = Inches(3 / 10)
    

table_shape.cell(1,0).text = "Kecelakaan"
cell2 = table_shape.cell(1,1).text_frame.paragraphs[0].add_run()
cell2.text = "Link poop op osdop o skdapo dkp aoskdpks dopask dposkd apsdo kokd"
cell2.hyperlink.address = "https://www.google.com"
table_shape.cell(2,0).text = "Lalu Lintas"
cell2 = table_shape.cell(2,1).text_frame.paragraphs[0].add_run()
cell2.text = "Link poop op osdop o skdapo dkp aoskdpks dopask dposkd apsdo kokd"
cell2.hyperlink.address = "https://www.google.com"
table_shape.cell(3,0).text = "Kebakaran"
table_shape.cell(3,1).text = "Kebakaran"

for cell in iter_cells(table_shape):
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(9)
            
table_shape.cell(0,0).text = "Event - Isu Positif"
table_shape.cell(0,0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
table_shape.cell(0,1).text = "Link"
table_shape.cell(0,1).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    
#============================ TABLE ATAS 2 =============================    


table2 = slide.placeholders[22]
x, y, cx, cy = Inches(2), Inches(2), Inches(4), Inches(1.5)
table_shape2 = table2.insert_table(rows=4, cols=2).table
table_shape2.columns[0].width = Inches(2.6)
table_shape2.columns[1].width = Inches(2.6)


table_shape2.cell(1,0).text = "Kecelakaan"
table_shape2.cell(1,1).text = "text"
table_shape2.cell(2,0).text = "Lalu Lintas"
table_shape2.cell(2,1).text = "link"
table_shape2.cell(3,0).text = "Kebakaran"
table_shape2.cell(3,1).text = "Kebakaran"

for cell in iter_cells(table_shape2):
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(9)
            
table_shape2.cell(0,0).text = "Event - Isu Positif"
table_shape2.cell(0,0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
table_shape2.cell(0,0).fill.solid()
table_shape2.cell(0,0).fill.fore_color.rgb = RGBColor(255,50,50)
table_shape2.cell(0,1).text = "Link"
table_shape2.cell(0,1).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
table_shape2.cell(0,1).fill.solid()
table_shape2.cell(0,1).fill.fore_color.rgb = RGBColor(255,50,50)

#================== TABLE BAWAH 1 =============================

table = slide.placeholders[20]
x, y, cx, cy = Inches(2), Inches(2), Inches(4), Inches(1.5)
table_shape = table.insert_table(rows=4, cols=2).table
table_shape.columns[0].width = Inches(2.6)
table_shape.columns[1].width = Inches(2.6)
for i in range(4):
    table_shape.rows[i].height = Inches(3 / 10)
    
table_shape.cell(0,0).text = "Event - Isu Positif"
table_shape.cell(0,0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
table_shape.cell(0,1).text = "Link"
table_shape.cell(0,1).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER



table_shape.cell(1,0).text = "Kecelakaan"
table_shape.cell(1,1).text = "text"
table_shape.cell(2,0).text = "Lalu Lintas"
table_shape.cell(2,1).text = "link"
table_shape.cell(3,0).text = "Kebakaran"
table_shape.cell(3,1).text = "Kebakaran"

table2 = slide.placeholders[23]
x, y, cx, cy = Inches(2), Inches(2), Inches(4), Inches(1.5)
table_shape2 = table2.insert_table(rows=4, cols=2).table
table_shape2.columns[0].width = Inches(2.6)
table_shape2.columns[1].width = Inches(2.6)


#============================ TABLE BAWAH 2 =============================
    
table_shape2.cell(0,0).text = "Event - Isu Positif"
table_shape2.cell(0,0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
table_shape2.cell(0,0).fill.solid()
table_shape2.cell(0,0).fill.fore_color.rgb = RGBColor(255,50,50)
table_shape2.cell(0,1).text = "Link"
table_shape2.cell(0,1).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
table_shape2.cell(0,1).fill.solid()
table_shape2.cell(0,1).fill.fore_color.rgb = RGBColor(255,50,50)

table_shape2.cell(1,0).text = "Kecelakaan"
table_shape2.cell(1,1).text = "text"
table_shape2.cell(2,0).text = "Lalu Lintas"
table_shape2.cell(2,1).text = "link"
table_shape2.cell(3,0).text = "Kebakaran"
table_shape2.cell(3,1).text = "Kebakaran"



slide3 = prs.slides.add_slide(prs.slide_layouts[2])
for shape in slide3.placeholders:
    print('%d %s' % (shape.placeholder_format.idx, shape.name))
 
 
date_slide3 = slide3.placeholders[10]
date_slide3.text = "19 November 2022"
title_slide3 = slide3.placeholders[11]
title_slide3.text = "POLRES ACEH BARAT DAYA - UNGGAHAN"


#============================ TABLE ATAS SLIDE 3 =============================
    
table_slide3 = slide3.placeholders[12]
x, y, cx, cy = Inches(2), Inches(2), Inches(4), Inches(1.5)
table_shape1_slide3 = table_slide3.insert_table(rows=11,cols=2).table
table_shape1_slide3.columns[0].width = Inches(3.5)
table_shape1_slide3.columns[1].width = Inches(3.5)

table_shape1_slide3.cell(0,0).text = "Event - Isu Positif"
table_shape1_slide3.cell(0,0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
table_shape1_slide3.cell(0,1).text = "Link"
table_shape1_slide3.cell(0,1).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER



table_slide3_2 = slide3.placeholders[13]
x, y, cx, cy = Inches(2), Inches(2), Inches(4), Inches(1.5)
table_shape2_slide3 = table_slide3_2.insert_table(rows=11,cols=2).table
table_shape2_slide3.columns[0].width = Inches(3.5)
table_shape2_slide3.columns[1].width = Inches(3.5)

table_shape2_slide3.cell(0,0).text = "Event - Isu Positif"
table_shape2_slide3.cell(0,0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
table_shape2_slide3.cell(0,0).fill.solid()
table_shape2_slide3.cell(0,0).fill.fore_color.rgb = RGBColor(255,50,50)
table_shape2_slide3.cell(0,1).text = "Link"
table_shape2_slide3.cell(0,1).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
table_shape2_slide3.cell(0,1).fill.solid()
table_shape2_slide3.cell(0,1).fill.fore_color.rgb = RGBColor(255,50,50)
       
                
prs.save('test.pptx')
os.startfile("test.pptx")