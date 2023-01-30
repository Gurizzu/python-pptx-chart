import collections 
import collections.abc
import json
import os
import traceback
from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.chart.axis import _BaseAxis
from pptx.enum.chart import XL_CATEGORY_TYPE
from pptx.dml.color import ColorFormat
from pptx.chart.series import _BaseSeries
from pptx.dml.chtfmt import ChartFormat
from pptx.dml.fill import FillFormat
from pptx.enum.chart import XL_TICK_MARK
from pptx.enum.chart import XL_TICK_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE

def ukuran(data):
        try:
            return round((data / 96),2)
        except Exception:
            traceback.print_exc()

with open(R"src\template_code\config.json","r") as f:
    data = json.load(f)
    
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


slides = data.get('slides')
for slide in slides:
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    for widget in slide.get("widgets"):
        match widget.get("widget").get("type"):
            case "chart":
                chart_data = ChartData()
                chart_data.categories = widget.get("widget").get("widget_data").get("categories")
                
                for chart_series in widget.get("widget").get("widget_data").get("series"):
                    chart_data.add_series(chart_series.get("legend"),chart_series.get("categories_point"))
                   
                                        
                x, y, cx, cy = Inches(widget.get("left")), Inches(widget.get("top")), Inches(widget.get("height")), Inches(widget.get("width"))
                
                chart_type = widget.get("widget").get("widget_ops").get("type")

                chart = slide1.shapes.add_chart(
                    XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
                ).chart
                
                if widget.get("widget").get("has_title") == True:
                    chart.has_title = True
                else:
                    chart.has_title = False
                
                if widget.get("widget").get("widget_ops").get("format").get("axis").get("data_label").get("has_data_labels") == True:
                    data_label = chart.plots[0]
                    data_label.has_data_labels = True
                    data_labels = data_label.data_labels
                    match widget.get("widget").get("widget_ops").get("format").get("axis").get("data_label").get("position"):
                        case "OUTSIDE_END":
                            data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
                        case "ABOVE":
                            data_labels.position = XL_LABEL_POSITION.ABOVE
                        case "BELOW":
                            data_labels.position = XL_LABEL_POSITION.BELOW                           
                    pt = widget.get("widget").get("widget_ops").get("format").get("axis").get("data_label").get("font_size")
                    data_labels.font.size = Pt(pt)
                
                value_axis = chart.value_axis
                axis_value = widget.get("widget").get("widget_ops").get("format").get("axis").get("value")
                
                if axis_value.get("visible") == False:
                    value_axis.visible = False
                else:
                    value_axis.visible = True
                    
                if axis_value.get("has_major_gridlines") == False:
                    value_axis.has_major_gridlines = False
                
                match axis_value.get("major_tick_mark"):
                    case "NONE":
                        value_axis.major_tick_mark = XL_TICK_MARK.NONE
                    case "INSIDE":
                        value_axis.major_tick_mark = XL_TICK_MARK.INSIDE
                        
                value_axis.tick_labels.number_format = axis_value.get("number_format")
                value_axis.tick_labels.font.size = Pt(axis_value.get("font_size"))
                value_axis.tick_labels.font.color.rgb = RGBColor.from_string(axis_value.get("color"))
                
                # match axis_value.get("tick_label_position"):
                #     case "HIGH":
                #         value_axis.major_tick_mark = XL_TICK_LABEL_POSITION.HIGH
                
                if axis_value.get("reverse_order") == False:                    
                    value_axis.reverse_order = False
                else:
                    value_axis.reverse_order = True
                    
                bar_plot = chart.plots[0]
                bar_plot.gap_width = widget.get("widget").get("widget_ops").get("format").get("chart_series").get("bar_plot_gap_width")
                   
                category_value = widget.get("widget").get("widget_ops").get("format").get("axis").get("category")   
                category_axis = chart.category_axis
                if category_value.get("reverse_order") == False:
                    category_axis.reverse_order = False
                else:
                    category_axis.reverse_order = True
                    
                match category_value.get("major_tick_mark"):
                    case "NONE":
                        value_axis.major_tick_mark = XL_TICK_MARK.NONE
                    case "INSIDE":
                        value_axis.major_tick_mark = XL_TICK_MARK.NONE
                category_axis.tick_labels.font.size = Pt(category_value.get("font_size"))
                
                if widget.get("widget").get("widget_ops").get("format").get("legend").get("has_legend") == False:
                    chart.has_legend = False
                else:
                    chart.has_legend = True
                    
                chart_series = chart.series[0]
                chart_series.format.fill.solid()
                chart_series.format.fill.fore_color.rgb = RGBColor.from_string(widget.get("widget").get("widget_ops").get("format").get("chart_series").get("color"))
            
            
            case "shape":
                left, top, width, height = Inches(widget.get("left")), Inches(widget.get("top")), Inches(widget.get("width")), Inches(widget.get("height"))
                
                shape = slide1.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, left, top, width, height
                )
                
                shape.fill.rgb = RGBColor.from_string("D84E2E")
                    
                 
                
                    
                
                

prs.save('chart_example5.pptx')                        
os.startfile("chart_example5.pptx")
